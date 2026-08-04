from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
ASSETS = ROOT / "report" / "assets"
OUT = ROOT / "ppt" / "DG_OXIDE_TFT_PDK_2026-08-06_v01.pptx"

NAVY = RGBColor(20, 35, 54)
INK = RGBColor(31, 42, 51)
MUTED = RGBColor(92, 105, 113)
TEAL = RGBColor(0, 128, 128)
RED = RGBColor(184, 56, 56)
AMBER = RGBColor(184, 119, 34)
PALE = RGBColor(241, 246, 247)
WHITE = RGBColor(255, 255, 255)


def text_box(slide, x, y, w, h, text, size=20, color=INK, bold=False, align=PP_ALIGN.LEFT):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.06)
    tf.margin_right = Inches(0.06)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = "Microsoft YaHei"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return shape


def bullet_box(slide, x, y, w, h, items, size=18, color=INK):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(0.08)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.name = "Microsoft YaHei"
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.space_after = Pt(9)
        p.bullet = True
    return shape


def base(slide, title, subtitle=None, number=None):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE
    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(0.18))
    band.fill.solid()
    band.fill.fore_color.rgb = TEAL
    band.line.fill.background()
    text_box(slide, 0.55, 0.35, 11.4, 0.55, title, 27, NAVY, True)
    if subtitle:
        text_box(slide, 0.58, 0.91, 11.4, 0.35, subtitle, 11, MUTED)
    if number is not None:
        text_box(slide, 12.35, 0.42, 0.45, 0.28, str(number), 12, TEAL, True, PP_ALIGN.RIGHT)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.58), Inches(7.16), Inches(12.18), Inches(0.01))
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(220, 228, 229)
    line.line.fill.background()


def card(slide, x, y, w, h, heading, body, accent=TEAL):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = PALE
    sh.line.color.rgb = RGBColor(210, 222, 224)
    sh.line.width = Pt(0.8)
    text_box(slide, x + 0.14, y + 0.10, w - 0.28, 0.34, heading, 16, accent, True)
    text_box(slide, x + 0.14, y + 0.52, w - 0.28, h - 0.62, body, 14, INK)


def image(slide, name, x, y, w=None, h=None):
    path = ASSETS / name
    if not path.exists():
        return
    slide.shapes.add_picture(str(path), Inches(x), Inches(y), width=Inches(w) if w else None, height=Inches(h) if h else None)


prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]

# 1
s = prs.slides.add_slide(blank)
s.background.fill.solid(); s.background.fill.fore_color.rgb = NAVY
text_box(s, 0.7, 1.0, 11.8, 0.7, "基于双栅 IGZO TFT 的二维器件模型与单极性逻辑教学 PDK", 30, WHITE, True)
text_box(s, 0.75, 1.9, 10.5, 0.5, "课程项目阶段性报告与可复现证据汇总", 22, RGBColor(182, 224, 224))
card(s, 0.75, 3.0, 3.65, 1.55, "范围", "仅 n 型 IGZO；二维优先；普通笔记本可运行", RGBColor(93, 202, 194))
card(s, 4.8, 3.0, 3.65, 1.55, "证据", "E2/E3 教学模型数值闭环；失败结果原样保留", RGBColor(93, 202, 194))
card(s, 8.85, 3.0, 3.65, 1.55, "当前门", "C00 R04 锚点未通过，G3 关闭，后续电路与版图关闭", RGBColor(239, 166, 117))
text_box(s, 0.75, 6.35, 11.7, 0.45, "提交版本：4d18f0b | 2026-08-04 | main 与 origin/main 同步", 13, RGBColor(190, 204, 210))

# 2
s = prs.slides.add_slide(blank); base(s, "研究问题与阶段边界", "把二维双栅器件证据推进到可审计的模型和逻辑评估", 2)
card(s, 0.7, 1.45, 3.8, 2.0, "器件层", "双栅 IGZO 漂移扩散与准静态陷阱耦合\n输出：Id-Vg、Id-Vd、状态和守恒", TEAL)
card(s, 4.78, 1.45, 3.8, 2.0, "模型层", "教学紧凑模型与 ngspice/GPL-Xyce 两路线\n输出：训练/验证误差与路线差异", TEAL)
card(s, 8.86, 1.45, 3.8, 2.0, "电路层", "双栅 IGZO 有源负载单极性反相器\n输出：VTC、延迟、功耗和逻辑门", RED)
text_box(s, 0.75, 4.1, 11.7, 0.35, "阶段门原则", 18, NAVY, True)
bullet_box(s, 0.8, 4.55, 11.6, 1.55, ["没有通过 C00，不能把环振、全加器或版图结果写成已验证。", "没有实验数据时，所有参数和趋势均限定为教学模型或文献约束敏感性。"], 18)

# 3
s = prs.slides.add_slide(blank); base(s, "双栅 IGZO 器件与模型边界", "二维结果为主，模型结果只在冻结教学域内解释", 3)
text_box(s, 0.75, 1.35, 5.5, 0.4, "结构与耦合", 19, NAVY, True)
bullet_box(s, 0.8, 1.85, 5.25, 2.25, ["n 型 IGZO；底栅/顶栅双介质耦合。", "二维漂移扩散、准静态陷阱电荷和电流守恒。", "Al2O3 物理厚度与 SPICE 有效 TOX 分开记录。"], 17)
text_box(s, 6.55, 1.35, 5.6, 0.4, "模型有效域", 19, NAVY, True)
card(s, 6.55, 1.85, 2.6, 1.35, "温度", "300 K 名义教学模型", TEAL)
card(s, 9.45, 1.85, 2.6, 1.35, "VDS", "0 至 0.2 V", TEAL)
card(s, 6.55, 3.45, 2.6, 1.35, "长度", "8 至 12 um", TEAL)
card(s, 9.45, 3.45, 2.6, 1.35, "材料", "IGZO only", TEAL)
text_box(s, 0.8, 5.45, 11.2, 0.65, "禁止外推：实验拟合、真实迁移率/DOS/接触参数、完整 P2/T03 物理结论。", 20, RED, True)

# 4
s = prs.slides.add_slide(blank); base(s, "双栅基线：从器件到耦合", "T02 关闭冻结教学模型的双向偏压与耦合数值门", 4)
image(s, "tcad_t02_c_bidirectional_families.png", 0.65, 1.35, w=7.25)
image(s, "tcad_t02_c_state_maps.png", 8.25, 1.35, w=4.45)
text_box(s, 0.75, 6.15, 11.8, 0.55, "T02-C：318 次 DC、248 点和 6 个状态；双向曲线、回程和耦合指标由独立检查复算。证据等级 E2（合同 E3）。", 15, MUTED)

# 5
s = prs.slides.add_slide(blank); base(s, "T03 敏感性结果", "五组参数均限定为冻结 IGZO 教学模型内的数值代理", 5)
image(s, "tcad_t03_p1_bias_sensitivity.png", 0.55, 1.25, w=4.05)
image(s, "tcad_t03_p2_bulk_traps_formal_v3_sensitivity.png", 4.65, 1.25, w=4.05)
image(s, "tcad_t03_p5_temperature_sensitivity.png", 8.75, 1.25, w=4.05)
text_box(s, 0.75, 6.2, 11.7, 0.55, "P1/P2/P3/P4/P5 正式运行与独立检查均已按阶段门闭环；方向是数值诊断，不是物理机制或实验校准。", 15, MUTED)

# 6
s = prs.slides.add_slide(blank); base(s, "M00/M01：紧凑模型与两路线对照", "结果可复现，但仍是行为等效教学候选", 6)
image(s, "m00_compact_model_fit_r02.png", 0.65, 1.35, w=5.6)
image(s, "m01_open_source_cross_check_r03.png", 6.55, 1.35, w=6.1)
card(s, 0.8, 5.55, 3.7, 0.9, "M00", "R02 runner E2 / 独立 E3；仅关闭冻结数值域", TEAL)
card(s, 4.82, 5.55, 3.7, 0.9, "M01", "ngspice + GPL Xyce；两路线教学候选", TEAL)
card(s, 8.84, 5.55, 3.7, 0.9, "边界", "不等同方程身份或真实器件参数", RED)

# 7
s = prs.slides.add_slide(blank); base(s, "C00 R04：完整执行但锚点失败", "这是一条真实的电路阶段门负结果，不是工具崩溃", 7)
image(s, "c00_active_load_inverter_vtc_r04.png", 0.55, 1.25, w=7.5)
card(s, 8.4, 1.45, 4.15, 1.2, "执行证据", "ngspice/Xyce 各 DC + 瞬态，共 4 个串行进程，均返回 0。", TEAL)
card(s, 8.4, 2.9, 4.15, 1.2, "锚点结果", "VOH 约 0.0912 V；VOL 约 0.0206 V；最大增益约 0.884。", RED)
card(s, 8.4, 4.35, 4.15, 1.2, "失败门", "0 个单位增益交点；无 VIL/VIH/NM；静态锚点不合格。", RED)
text_box(s, 0.75, 6.25, 11.8, 0.42, "33/36 E0/FAIL；所有原始输出、表、差异、图和哈希均保留。", 16, RED, True)

# 8
s = prs.slides.add_slide(blank); base(s, "C00 瞬态与 G3 阶段门", "没有把缺失 crossing 伪装成延迟", 8)
image(s, "c00_active_load_inverter_transient_r04.png", 0.55, 1.2, w=8.0)
bullet_box(s, 8.75, 1.45, 3.65, 2.7, ["输出约 0.022–0.057 V。", "rise/fall crossing 均不存在。", "tPHL/tPLH 保留为空，不插值、不估算。", "因此 C00 未通过，C01/C02/C03、版图和 PEX 关闭。"], 17)
text_box(s, 8.78, 5.2, 3.45, 0.8, "结论：设计点不满足预注册逻辑门，而不是仿真器不收敛。", 20, RED, True)

# 9
s = prs.slides.add_slide(blank); base(s, "交付边界与风险控制", "当前版本适合课程提交，不能冒充流片级 PDK", 9)
card(s, 0.75, 1.35, 3.8, 2.0, "已完成", "二维器件数值\nT03 五组敏感性\nM00/M01 教学模型\n完整证据矩阵", TEAL)
card(s, 4.8, 1.35, 3.8, 2.0, "明确未完成", "合格 C00 反相器\n组合逻辑与环振\n版图/LVS/PEX\nHZO 铁电扩展", RED)
card(s, 8.85, 1.35, 3.8, 2.0, "不能夸大", "实验校准\n物理参数提取\n工艺签核\n真实工作频率与可靠性", AMBER)
text_box(s, 0.8, 4.15, 11.7, 0.42, "阶段门策略把失败保留下来，避免用后续占位结果掩盖当前电路不可用。", 20, NAVY, True)
bullet_box(s, 0.85, 4.75, 11.2, 1.3, ["下一轮 R05 必须重新预注册设计变量与锚点，再独立运行；不能回写 R04。", "本次提交版本、报告和哈希链可从 Git 提交 4d18f0b 复现。"], 17)

# 10
s = prs.slides.add_slide(blank); base(s, "最终结论", "以可复现、边界清晰的结果完成课程交付", 10)
text_box(s, 0.85, 1.45, 11.5, 0.55, "双栅 IGZO 二维教学 PDK 的器件与模型主线已闭环；有源负载逻辑在冻结锚点上被数据判定为不合格。", 24, NAVY, True)
card(s, 0.85, 2.45, 3.75, 1.7, "贡献 1", "建立二维双栅 IGZO 器件、陷阱与敏感性证据链。", TEAL)
card(s, 4.8, 2.45, 3.75, 1.7, "贡献 2", "建立可复核的行为紧凑模型和开源两路线对照。", TEAL)
card(s, 8.75, 2.45, 3.75, 1.7, "贡献 3", "用正式 C00 失败暴露设计裕量不足并关闭下游门。", RED)
text_box(s, 0.9, 5.05, 11.2, 0.6, "证据等级：器件/模型主线 E2-E3；C00 R04 runner E0/FAIL。", 21, INK, True)
text_box(s, 0.9, 5.75, 11.2, 0.45, "版本：4d18f0b | 仓库与报告检查通过 | IGZO only", 15, MUTED)

OUT.parent.mkdir(parents=True, exist_ok=True)
prs.save(OUT)
print(OUT)
